# """
# dua_core.py — pure string replacement, zero minidom serialisation.
# Preserves original zip compression type for every file.
# """

# import io, zipfile, re, html


# # ── parser ────────────────────────────────────────────────────────────────────

# def parse_duas(text: str) -> list:
#     sets, buf = [], []
#     for raw in text.splitlines():
#         s = raw.strip()
#         if s:
#             buf.append(s)
#         else:
#             if len(buf) >= 4:
#                 sets.append({"arabic": buf[0], "transliteration": buf[1],
#                               "english": buf[2], "urdu": buf[3]})
#             buf = []
#     if len(buf) >= 4:
#         sets.append({"arabic": buf[0], "transliteration": buf[1],
#                      "english": buf[2], "urdu": buf[3]})
#     return sets


# def _esc(text: str) -> str:
#     return html.escape(text, quote=False)


# # ── per-shape string replacement ──────────────────────────────────────────────

# def _replace_shape_text(xml: str, shape_id: str, new_text: str) -> str:
#     escaped = _esc(new_text)
#     idx = xml.find(f'id="{shape_id}"')
#     if idx == -1:
#         return xml
#     sp_end = xml.find("</p:sp>", idx) + len("</p:sp>")
#     block = xml[idx:sp_end]
#     block = re.sub(r'<a:t>[^<]*</a:t>', f'<a:t>{escaped}</a:t>', block, count=1)
#     return xml[:idx] + block + xml[sp_end:]


# def _replace_transliteration(xml: str, shape_id: str, new_text: str) -> str:
#     escaped = _esc(new_text)
#     idx = xml.find(f'id="{shape_id}"')
#     if idx == -1:
#         return xml
#     sp_end = xml.find("</p:sp>", idx) + len("</p:sp>")
#     block = xml[idx:sp_end]
#     ap_match = re.search(r'<a:p>(.*?)</a:p>', block, re.DOTALL)
#     if not ap_match:
#         return xml
#     para_inner = ap_match.group(1)
#     runs = re.findall(r'<a:r>.*?</a:r>', para_inner, re.DOTALL)
#     if not runs:
#         return xml
#     first_run = re.sub(r'<a:t>[^<]*</a:t>', f'<a:t>{escaped}</a:t>', runs[0], count=1)
#     pPr_match = re.search(r'<a:pPr[^>]*/?>(?:</a:pPr>)?', para_inner)
#     pPr = pPr_match.group(0) if pPr_match else ''
#     new_para_inner = pPr + first_run
#     new_block = block[:ap_match.start(1)] + new_para_inner + block[ap_match.end(1):]
#     return xml[:idx] + new_block + xml[sp_end:]


# def _fix_arabic_and_english(xml: str) -> str:
#     # Shape 106 (Arabic): remove placeholder, centre-align
#     xml = xml.replace('<p:ph type="ctrTitle"/>', '')
#     xml = re.sub(
#         r'(<p:cNvPr id="106"[^/]*/><p:cNvSpPr).*?</p:cNvSpPr>',
#         r'\1 txBox="1"></p:cNvSpPr>',
#         xml, flags=re.DOTALL, count=1)
#     idx = xml.find('id="106"')
#     if idx != -1:
#         sp_end = xml.find("</p:sp>", idx) + len("</p:sp>")
#         block = xml[idx:sp_end]
#         block = re.sub(r'(<a:pPr\b)([^>]*?)(/>|>)',
#                        lambda m: m.group(1) + m.group(2) +
#                        (' algn="ctr"' if 'algn' not in m.group(2) else '') +
#                        m.group(3),
#                        block, count=1)
#         xml = xml[:idx] + block + xml[sp_end:]

#     # Shape 107 (English): remove placeholder, fix style
#     xml = re.sub(r'<p:ph type="subTitle"[^/]*/>', '', xml)
#     idx = xml.find('id="107"')
#     if idx != -1:
#         sp_end = xml.find("</p:sp>", idx) + len("</p:sp>")
#         block = xml[idx:sp_end]
#         block = re.sub(
#             r'(<p:cNvPr id="107"[^/]*/><p:cNvSpPr).*?</p:cNvSpPr>',
#             r'\1 txBox="1"></p:cNvSpPr>',
#             block, flags=re.DOTALL, count=1)
#         block = re.sub(r'<a:pPr[^>]*>', '<a:pPr algn="ctr">', block, count=1)
#         block = re.sub(
#             r'<a:rPr[^>]*>.*?(?=<a:t>)',
#             '<a:rPr lang="en-US" sz="3200" dirty="0"><a:solidFill>'
#             '<a:srgbClr val="0066CC"/></a:solidFill>',
#             block, count=1, flags=re.DOTALL)
#         xml = xml[:idx] + block + xml[sp_end:]

#     return xml


# def _fill_slide(xml_bytes: bytes, dua: dict) -> bytes:
#     xml = xml_bytes.decode("utf-8")
#     xml = _fix_arabic_and_english(xml)
#     xml = _replace_shape_text(xml, "106", dua["arabic"])
#     xml = _replace_shape_text(xml, "107", dua["english"])
#     xml = _replace_shape_text(xml, "109", dua["urdu"])
#     xml = _replace_transliteration(xml, "108", dua["transliteration"])
#     return xml.encode("utf-8")


# # ── main public function ──────────────────────────────────────────────────────

# def build_pptx_bytes(template_bytes: bytes, duas: list) -> bytes:
#     with zipfile.ZipFile(io.BytesIO(template_bytes), "r") as zin:
#         infos     = {i.filename: i for i in zin.infolist()}
#         file_data = {n: zin.read(n) for n in infos}

#     slide1_xml  = file_data["ppt/slides/slide1.xml"]
#     slide1_rels = file_data.get("ppt/slides/_rels/slide1.xml.rels", b"")

#     prs_xml           = file_data["ppt/presentation.xml"].decode("utf-8")
#     prs_rels_xml      = file_data["ppt/_rels/presentation.xml.rels"].decode("utf-8")
#     content_types_xml = file_data["[Content_Types].xml"].decode("utf-8")

#     existing_ids  = [int(x) for x in re.findall(r'<p:sldId[^>]+id="(\d+)"', prs_xml)]
#     next_slide_id = max(existing_ids) + 1 if existing_ids else 256
#     existing_rids = [int(x) for x in re.findall(r'Id="rId(\d+)"', prs_rels_xml)]
#     next_rid_num  = max(existing_rids) + 1 if existing_rids else 1
#     existing_nums = [int(m.group(1)) for n in infos
#                      if (m := re.match(r"ppt/slides/slide(\d+)\.xml", n))]
#     next_slide_num = max(existing_nums) + 1 if existing_nums else 2

#     new_files = {}

#     for i, dua in enumerate(duas):
#         if i == 0:
#             new_files["ppt/slides/slide1.xml"] = _fill_slide(slide1_xml, dua)
#         else:
#             new_num   = next_slide_num + (i - 1)
#             slide_key = f"ppt/slides/slide{new_num}.xml"
#             rels_key  = f"ppt/slides/_rels/slide{new_num}.xml.rels"

#             new_files[slide_key] = _fill_slide(slide1_xml, dua)

#             if slide1_rels:
#                 rels_str = re.sub(r'\s*<Relationship[^>]*notesSlide[^>]*/>\s*',
#                                   "\n", slide1_rels.decode("utf-8"))
#                 new_files[rels_key] = rels_str.encode("utf-8")

#             ct_entry = (f'<Override PartName="/ppt/slides/slide{new_num}.xml" '
#                         f'ContentType="application/vnd.openxmlformats-officedocument'
#                         f'.presentationml.slide+xml"/>')
#             if f"/ppt/slides/slide{new_num}.xml" not in content_types_xml:
#                 content_types_xml = content_types_xml.replace(
#                     "</Types>", f"  {ct_entry}\n</Types>")

#             rid = f"rId{next_rid_num + (i - 1)}"
#             new_rel = (f'<Relationship Id="{rid}" '
#                        f'Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide" '
#                        f'Target="slides/slide{new_num}.xml"/>')
#             prs_rels_xml = prs_rels_xml.replace(
#                 "</Relationships>", f"  {new_rel}\n</Relationships>")

#             sld_id = next_slide_id + (i - 1)

#             # Update main sldIdLst
#             new_sld = f'<p:sldId id="{sld_id}" r:id="{rid}"/>'
#             prs_xml = prs_xml.replace("</p:sldIdLst>",
#                                       f"{new_sld}</p:sldIdLst>")

#             # Also update the p14:sldIdLst (section tracking) — must stay in sync
#             new_p14_sld = f'<p14:sldId id="{sld_id}"/>'
#             prs_xml = prs_xml.replace("</p14:sldIdLst>",
#                                       f"{new_p14_sld}</p14:sldIdLst>")

#     new_files["ppt/presentation.xml"]            = prs_xml.encode("utf-8")
#     new_files["ppt/_rels/presentation.xml.rels"] = prs_rels_xml.encode("utf-8")
#     new_files["[Content_Types].xml"]             = content_types_xml.encode("utf-8")

#     out = io.BytesIO()
#     with zipfile.ZipFile(out, "w", zipfile.ZIP_DEFLATED) as zout:
#         for name in infos:
#             data     = new_files.get(name, file_data[name])
#             compress = infos[name].compress_type
#             zout.writestr(zipfile.ZipInfo(name), data, compress_type=compress)
#         for name, data in new_files.items():
#             if name not in infos:
#                 zout.writestr(name, data, compress_type=zipfile.ZIP_DEFLATED)

#     return out.getvalue()
"""
dua_core.py — uses python-pptx for reliable PPTX generation.
Duplicates slides by copying the XML element directly into a new slide part.
"""

import io, copy
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.parts.slide import SlidePart
from lxml import etree


# ── parser ────────────────────────────────────────────────────────────────────

def parse_duas(text: str) -> list:
    sets, buf = [], []
    for raw in text.splitlines():
        s = raw.strip()
        if s:
            buf.append(s)
        else:
            if len(buf) >= 4:
                sets.append({"arabic": buf[0], "transliteration": buf[1],
                              "english": buf[2], "urdu": buf[3]})
            buf = []
    if len(buf) >= 4:
        sets.append({"arabic": buf[0], "transliteration": buf[1],
                     "english": buf[2], "urdu": buf[3]})
    return sets


# ── shape text setter ─────────────────────────────────────────────────────────

def _set_shape_text(shape, text, centre=False, font_size=None, colour=None):
    tf   = shape.text_frame
    para = tf.paragraphs[0]

    # Save rPr from first run to preserve font face / language
    ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    saved_rPr = None
    if para.runs:
        rPr_elem = para.runs[0]._r.find(f'{{{ns}}}rPr')
        if rPr_elem is not None:
            saved_rPr = copy.deepcopy(rPr_elem)

    # Remove extra paragraphs, clear runs from first para
    for p in tf.paragraphs[1:]:
        p._p.getparent().remove(p._p)
    for run in para.runs:
        para._p.remove(run._r)

    if centre:
        para.alignment = PP_ALIGN.CENTER

    run = para.add_run()
    run.text = text

    # Re-apply saved rPr (preserves Arabic/Urdu font, RTL, etc.)
    if saved_rPr is not None:
        existing = run._r.find(f'{{{ns}}}rPr')
        if existing is not None:
            run._r.remove(existing)
        new_rPr = copy.deepcopy(saved_rPr)

        # Override font size
        if font_size is not None:
            new_rPr.set('sz', str(int(font_size * 100)))

        # Override colour
        if colour is not None:
            for sf in new_rPr.findall(f'{{{ns}}}solidFill'):
                new_rPr.remove(sf)
            hex_col = ''.join(f'{c:02X}' for c in colour)
            sf_elem = etree.fromstring(
                f'<a:solidFill xmlns:a="{ns}">'
                f'<a:srgbClr val="{hex_col}"/></a:solidFill>')
            new_rPr.insert(0, sf_elem)

        run._r.insert(0, new_rPr)
    else:
        if font_size is not None:
            run.font.size = Pt(font_size)
        if colour is not None:
            run.font.color.rgb = RGBColor(*colour)


def _fill_slide(slide, dua: dict):
    for shape in slide.shapes:
        sid = shape.shape_id
        if not shape.has_text_frame:
            continue
        if sid == 106:    # Arabic
            _set_shape_text(shape, dua["arabic"], centre=True)
        elif sid == 107:  # English
            _set_shape_text(shape, dua["english"], centre=True,
                            font_size=32, colour=(0x00, 0x66, 0xCC))
        elif sid == 109:  # Urdu
            _set_shape_text(shape, dua["urdu"])
        elif sid == 108:  # Transliteration
            _set_shape_text(shape, dua["transliteration"], centre=True,
                            font_size=26.4, colour=(0x00, 0x66, 0xCC))


def _duplicate_slide(prs, source_slide):
    """Add a new slide that is a deep copy of source_slide."""
    template   = source_slide
    slide_layout = template.slide_layout

    # Create a new blank slide part
    prs_part   = prs.part
    partname   = prs_part._next_slide_partname
    slide_part = SlidePart.new(partname, prs_part.package, slide_layout.part)

    # Replace its XML with a deep copy of the template slide's XML
    slide_part._element = copy.deepcopy(template._element)

    # Register relationship and slide ID
    rId = prs_part.relate_to(slide_part, RT.SLIDE)
    prs.slides._sldIdLst.add_sldId(rId)

    return slide_part.slide


# ── main public function ──────────────────────────────────────────────────────

def build_pptx_bytes(template_bytes: bytes, duas: list) -> bytes:
    prs            = Presentation(io.BytesIO(template_bytes))
    template_slide = prs.slides[0]

    # Fill slide 1
    _fill_slide(template_slide, duas[0])

    # Duplicate and fill for remaining duas
    for dua in duas[1:]:
        new_slide = _duplicate_slide(prs, template_slide)
        _fill_slide(new_slide, dua)

    out = io.BytesIO()
    prs.save(out)
    return out.getvalue()